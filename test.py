from pptx import Presentation
from pptx.util import Inches
from math import cos, sin, radians
from pptx.dml.color import RGBColor


class RoundedRectangle:
    def __init__(self, slide, left, top, width, height, corner_radius, steps=10):
        """
        Initialize the RoundedRectangle object.
        :param slide: PowerPoint slide object where the shape will be added.
        :param left: x-coordinate of the top-left corner of the rectangle.
        :param top: y-coordinate of the top-left corner of the rectangle.
        :param width: Width of the rectangle.
        :param height: Height of the rectangle.
        :param corner_radius: Radius of the rounded corners.
        :param steps: Number of steps to approximate the rounded corner (higher = smoother).
        """
        self.slide = slide
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.corner_radius = corner_radius
        self.steps = steps
        self.shape = None  # The generated shape object

    def _add_rounded_corner(self, builder, center_x, center_y, start_angle, end_angle, radius, steps):
        """
        Add a rounded corner to the builder.
        :param builder: FreeformBuilder object used to construct the shape.
        :param center_x: x-coordinate of the corner's center.
        :param center_y: y-coordinate of the corner's center.
        :param start_angle: Starting angle of the arc (in degrees).
        :param end_angle: Ending angle of the arc (in degrees).
        :param radius: Radius of the arc.
        :param steps: Number of steps to approximate the arc.
        """
        for step in range(steps + 1):
            angle = radians(start_angle + (end_angle - start_angle) * step / steps)
            x = center_x + radius * cos(angle)
            y = center_y + radius * sin(angle)
            builder.add_line_segments([(x, y)], close=False)

    def create(self):
        """
        Create the rounded rectangle shape and return it.
        """
        builder = self.slide.shapes.build_freeform(self.left + self.corner_radius, self.top)

        # Top edge + top-right corner
        builder.add_line_segments([(self.left + self.width - self.corner_radius, self.top)], close=False)
        self._add_rounded_corner(builder, self.left + self.width - self.corner_radius, self.top + self.corner_radius, 270, 360, self.corner_radius, self.steps)

        # Right edge + bottom-right corner
        builder.add_line_segments([(self.left + self.width, self.top + self.height - self.corner_radius)], close=False)
        self._add_rounded_corner(builder, self.left + self.width - self.corner_radius, self.top + self.height - self.corner_radius, 0, 90, self.corner_radius, self.steps)

        # Bottom edge + bottom-left corner
        builder.add_line_segments([(self.left + self.corner_radius, self.top + self.height)], close=False)
        self._add_rounded_corner(builder, self.left + self.corner_radius, self.top + self.height - self.corner_radius, 90, 180, self.corner_radius, self.steps)

        # Left edge + top-left corner
        builder.add_line_segments([(self.left, self.top + self.corner_radius)], close=False)
        self._add_rounded_corner(builder, self.left + self.corner_radius, self.top + self.corner_radius, 180, 270, self.corner_radius, self.steps)

        # Close the shape
        self.shape = builder.convert_to_shape()
        return self.shape

    def set_style(self, fill_color=None, line_color=None, line_width=None):
        """
        Set the style for the shape (fill color, line color, line width).
        :param fill_color: Fill color as an RGBColor object.
        :param line_color: Line color as an RGBColor object.
        :param line_width: Line width (in EMU).
        """
        if self.shape is None:
            raise ValueError("Shape has not been created. Call `create` first.")
        
        # Set fill color
        if fill_color:
            self.shape.fill.solid()
            self.shape.fill.fore_color.rgb = fill_color
        
        # Set line color and width
        if line_color:
            self.shape.line.color.rgb = line_color
        if line_width:
            self.shape.line.width = line_width

# Create a presentation and add a slide
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Create a RoundedRectangle
rounded_rect = RoundedRectangle(
    slide=slide,
    left=Inches(1),
    top=Inches(1),
    width=Inches(4),
    height=Inches(2),
    corner_radius=Inches(0.5),
    steps=20  # Higher steps create a smoother corner
)

# Create the shape
shape = rounded_rect.create()

# Apply style to the shape
rounded_rect.set_style(
    fill_color=RGBColor(255, 0, 0),    # Red fill
    line_color=RGBColor(0, 0, 0),     # Black border
    line_width=Inches(0.05)           # Border width
)
shape.text = "Rounded Rectangle"

# Save the presentation
prs.save("rounded_rectangle_class.pptx")
