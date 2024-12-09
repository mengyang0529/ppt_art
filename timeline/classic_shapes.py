from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from core.geometry import *
from core.roundedrectangle import RoundedRectangle
        
class MainTopicBox(Box):
    def __init__(self, slide, node, box):
        super().__init__(box.left, box.top, box.width, box.height)

        node.color = RGBColor(0, 0, 0)
        # Create a RoundedRectangle
        rounded_rect = RoundedRectangle(
            slide=slide,
            left=Inches(box.left),
            top=Inches(box.top),
            width=Inches(box.width),
            height=Inches(box.height),
            corner_radius=Inches(0.3),
            steps=20  # Higher steps create a smoother corner
        )
        # Create the shape
        shape = rounded_rect.create()
        self.convert_to_box(shape)
        # self.add_to_slide(slide, MSO_SHAPE.ROUNDED_RECTANGLE)
        self.text(Pt(16), RGBColor(255, 255, 255), node.value)
        self.box.fill.solid()
        self.box.fill.fore_color.rgb = node.color
        
        self.id = node.id

    def text(self, size, color, text):
        self.box.text = text
        self.box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.box.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.MIDDLE
        self.box.text_frame.paragraphs[0].font.bold = True
        self.box.text_frame.paragraphs[0].font.size = size
        self.box.text_frame.paragraphs[0].font.color.rgb = color

    def add_to_slide(self, slide, shape_type):
        self.box = slide.shapes.add_shape(
            shape_type,
            Inches(self.left), Inches(self.top), Inches(self.width), Inches(self.height)
        )

class SubTopicBox(Box):
    def __init__(self, slide, node, box) -> None:
        super().__init__(box.left, box.top, box.width, box.height)
        # Create a RoundedRectangle
        rounded_rect = RoundedRectangle(
            slide=slide,
            left=Inches(box.left),
            top=Inches(box.top),
            width=Inches(box.width),
            height=Inches(box.height),
            corner_radius=Inches(0.2),
            steps=40  # Higher steps create a smoother corner
        )
        self.colors = [
            RGBColor(238, 99, 99),
            RGBColor(244, 164, 96),
            RGBColor(255, 215, 0),
            RGBColor(60, 179, 113),
            RGBColor(70, 130, 180),
            RGBColor(106, 90, 205)
        ]
        # Create the shap
        shape = rounded_rect.create()
        self.convert_to_box(shape)
        # self.add_to_slide(slide, MSO_SHAPE.ROUNDED_RECTANGLE)
        self.text(Pt(14), RGBColor(0, 0, 0), node.value)
        self.box.fill.solid()
        node.color = self.colors[node.color_id % (len(self.colors))]
        self.box.line.color.rgb = node.color
        self.box.fill.fore_color.rgb = node.color
        self.box.line.width = Pt(2)
        self.id = node.id

    def text(self, size, color, text):
        self.box.text = text
        self.box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.box.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.MIDDLE
        self.box.text_frame.paragraphs[0].font.bold = True
        self.box.text_frame.paragraphs[0].font.size = size
        self.box.text_frame.paragraphs[0].font.color.rgb = color

    def add_to_slide(self, slide, shape_type):
        self.box = slide.shapes.add_shape(
            shape_type,
            Inches(self.left), Inches(self.top), Inches(self.width), Inches(self.height)
        )

class DetailBox(Box):
    def __init__(self, slide, node, box) -> None:
        super().__init__(box.left, box.top, box.width, box.height)

        self.add_to_slide(slide, MSO_SHAPE.RECTANGLE)
        self.text(Pt(12), RGBColor(0, 0, 0), node.value)
        self.box.fill.background()
        node.color = node.parent.color
        self.box.line.width = Pt(2)

        self.box.line.color.rgb = node.color
        self.id = node.id

    def text(self, size, color, text):
        self.box.text = text
        self.box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        self.box.text_frame.paragraphs[0].vertical_anchor = MSO_ANCHOR.MIDDLE
        self.box.text_frame.paragraphs[0].font.bold = True
        self.box.text_frame.paragraphs[0].font.size = size
        self.box.text_frame.paragraphs[0].font.color.rgb = color

    def add_to_slide(self, slide, shape_type):
        self.box = slide.shapes.add_shape(
            shape_type,
            Inches(self.left), Inches(self.top), Inches(self.width), Inches(self.height)
        )
    
class DrawingBox(Box):
    def __init__(self, slide, node, box, box_type) -> None:
        super().__init__(box.left, box.top, box.width, box.height)
        if box_type == 0:
            MainTopicBox(slide, node, box)
        elif box_type == 1:
            SubTopicBox(slide, node, box)
        elif box_type >= 2:
            DetailBox(slide, node, box)

class DrawingLine:
    def __init__(self, slide, node, pt1, pt2) -> None:

            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, 
                pt1.x, pt1.y, 
                pt2.x, pt2.y 
            )
            connector.line.color.rgb = node.color
            connector.line.width = Pt(2)

if __name__ == "__main__":
    pass
